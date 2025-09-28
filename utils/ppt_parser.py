"""
PowerPoint parser module for extracting text from PPT/PPTX pitch decks.

This module provides functionality to extract text content from PowerPoint
presentations using python-pptx, handling slides and text shapes.
"""

from pptx import Presentation


def parse_ppt(filepath):
    """
    Extract text content from a PowerPoint presentation file.
    
    This function reads a PPT or PPTX file and extracts all text content
    from text shapes across all slides, concatenating them into a single string.
    
    Args:
        filepath (str): Path to the PowerPoint file to parse (.ppt or .pptx).
        
    Returns:
        dict: Dictionary containing:
            - source (str): File format identifier ("ppt")
            - raw_text (str): Extracted text content from all slides
            
    Note:
        The function iterates through all slides and all shapes within each slide,
        extracting text only from shapes that have a 'text' attribute (text boxes,
        placeholders, etc.). Each shape's text is separated by a newline character.
    """
    prs = Presentation(filepath)
    text = ""
    
    # Extract text from all slides
    for slide in prs.slides:
        # Extract text from all shapes in the slide
        for shape in slide.shapes:
            # Only process shapes that have text content
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    
    return {
        "source": "ppt",
        "raw_text": text
    }
